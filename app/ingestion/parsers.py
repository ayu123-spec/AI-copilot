"""Parsers that turn raw files into a normalized ParsedDocument.

Each parser preserves page/slide structure so downstream chunks can cite a page.
Heavy/optional libraries are imported lazily inside each function so importing
this module never fails if a format's dependency is absent.
"""

from pathlib import Path

from app.ingestion.base import ParsedDocument, ParsedPage

# Map common extensions to a canonical content type.
EXTENSION_CONTENT_TYPE = {
    ".pdf": "application/pdf",
    ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ".pptx": (
        "application/vnd.openxmlformats-officedocument." "presentationml.presentation"
    ),
    ".txt": "text/plain",
    ".md": "text/markdown",
    ".markdown": "text/markdown",
}


class UnsupportedFileType(Exception):
    pass


def _parse_pdf(path: Path) -> ParsedDocument:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    pages = [
        ParsedPage(page_number=i + 1, text=(page.extract_text() or "").strip())
        for i, page in enumerate(reader.pages)
    ]
    meta = reader.metadata or {}
    metadata = {
        "author": getattr(meta, "author", None),
        "title": getattr(meta, "title", None),
        "source": path.name,
    }
    return ParsedDocument(path.name, "application/pdf", pages, metadata)


def _parse_docx(path: Path) -> ParsedDocument:
    from docx import Document as Docx

    doc = Docx(str(path))
    text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    props = doc.core_properties
    metadata = {"author": props.author, "title": props.title, "source": path.name}
    # python-docx has no real page concept; treat the body as a single page.
    return ParsedDocument(
        path.name,
        EXTENSION_CONTENT_TYPE[".docx"],
        [ParsedPage(page_number=1, text=text)],
        metadata,
    )


def _parse_pptx(path: Path) -> ParsedDocument:
    from pptx import Presentation

    prs = Presentation(str(path))
    pages = []
    for i, slide in enumerate(prs.slides):
        parts = [
            shape.text
            for shape in slide.shapes
            if shape.has_text_frame and shape.text.strip()
        ]
        pages.append(ParsedPage(page_number=i + 1, text="\n".join(parts)))
    return ParsedDocument(
        path.name, EXTENSION_CONTENT_TYPE[".pptx"], pages, {"source": path.name}
    )


def _parse_text(path: Path, content_type: str) -> ParsedDocument:
    text = path.read_text(encoding="utf-8", errors="replace")
    return ParsedDocument(
        path.name,
        content_type,
        [ParsedPage(page_number=1, text=text)],
        {"source": path.name},
    )


def parse_file(path: str | Path) -> ParsedDocument:
    """Dispatch to the correct parser based on file extension."""
    path = Path(path)
    ext = path.suffix.lower()
    if ext == ".pdf":
        return _parse_pdf(path)
    if ext == ".docx":
        return _parse_docx(path)
    if ext == ".pptx":
        return _parse_pptx(path)
    if ext in (".txt",):
        return _parse_text(path, "text/plain")
    if ext in (".md", ".markdown"):
        return _parse_text(path, "text/markdown")
    raise UnsupportedFileType(f"Unsupported file type: {ext}")
