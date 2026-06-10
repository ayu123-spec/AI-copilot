from pathlib import Path

import pytest

from app.ingestion.parsers import UnsupportedFileType, parse_file
from app.preprocessing.cleaner import clean_text


def test_parse_txt(tmp_path: Path):
    f = tmp_path / "note.txt"
    f.write_text("Hello world.\nSecond line.")
    doc = parse_file(f)
    assert doc.content_type == "text/plain"
    assert "Hello world" in doc.full_text
    assert doc.pages[0].page_number == 1


def test_parse_markdown(tmp_path: Path):
    f = tmp_path / "readme.md"
    f.write_text("# Title\n\nSome **bold** content.")
    doc = parse_file(f)
    assert doc.content_type == "text/markdown"
    assert "content" in doc.full_text


def test_parse_docx(tmp_path: Path):
    from docx import Document as Docx

    f = tmp_path / "memo.docx"
    d = Docx()
    d.add_paragraph("Quarterly results were strong.")
    d.add_paragraph("Revenue grew across regions.")
    d.save(f)

    doc = parse_file(f)
    assert "Quarterly results" in doc.full_text
    assert "Revenue grew" in doc.full_text


def test_parse_pptx(tmp_path: Path):
    from pptx import Presentation

    f = tmp_path / "deck.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Roadmap"
    prs.save(f)

    doc = parse_file(f)
    assert "Roadmap" in doc.full_text
    assert doc.pages[0].page_number == 1


def test_unsupported_type(tmp_path: Path):
    f = tmp_path / "data.xyz"
    f.write_text("nope")
    with pytest.raises(UnsupportedFileType):
        parse_file(f)


def test_clean_text_removes_page_numbers_and_whitespace():
    raw = "Heading\n\n\n\nPage 3 of 10\nReal    content   here\n12\n"
    cleaned = clean_text(raw)
    assert "Page 3 of 10" not in cleaned
    assert "\n12\n" not in cleaned
    assert "Real content here" in cleaned
